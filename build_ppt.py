"""Build the 'Achter de schermen' Wifipool chatbot PPT.

Output: Wifipool_Chatbot_Achter_de_Schermen.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand palette (matches dashboard ocean+primary) ──────────────────────
NAVY      = RGBColor(0x0B, 0x1F, 0x3A)
OCEAN_700 = RGBColor(0x0E, 0x5A, 0x90)
OCEAN_500 = RGBColor(0x0E, 0x7A, 0xC4)
OCEAN_400 = RGBColor(0x3E, 0xA5, 0xD8)
OCEAN_100 = RGBColor(0xD8, 0xF1, 0xFA)
OCEAN_50  = RGBColor(0xEE, 0xF8, 0xFD)
PRIMARY   = RGBColor(0x6C, 0x63, 0xFF)
PRIMARY_2 = RGBColor(0x8E, 0x86, 0xFF)
ACCENT    = RGBColor(0xFF, 0xC1, 0x07)
INK       = RGBColor(0x10, 0x1A, 0x2C)
INK_SOFT  = RGBColor(0x42, 0x52, 0x6B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE    = RGBColor(0xF3, 0xF6, 0xFC)
LINE      = RGBColor(0xE2, 0xE8, 0xF2)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
RED       = RGBColor(0xEF, 0x44, 0x44)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill, line=None, radius=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    if radius is not None:
        # Adjust corner radius (0.0–1.0)
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def add_oval(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri", line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_emoji(slide, x, y, size_emu, glyph, color=PRIMARY, bg=None):
    if bg is not None:
        add_oval(slide, x, y, size_emu, size_emu, bg)
        # center the glyph inside
        add_text(slide, x, y, size_emu, size_emu, glyph,
                 size=int(size_emu / 12700 / 2.4), color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Segoe UI Emoji")
    else:
        add_text(slide, x, y, size_emu, size_emu, glyph,
                 size=int(size_emu / 12700 / 2), color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Segoe UI Emoji")


def page_chrome(slide, kicker=None, title=None, page_num=None, total=None):
    """Adds a top accent bar, kicker, title, and page footer."""
    # Top accent
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.18), PRIMARY)
    # Logo dot
    add_oval(slide, Inches(0.45), Inches(0.42), Inches(0.32), Inches(0.32), OCEAN_500)
    add_text(slide, Inches(0.45), Inches(0.42), Inches(0.32), Inches(0.32),
             "🌊", size=14, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")
    add_text(slide, Inches(0.85), Inches(0.42), Inches(3.5), Inches(0.32),
             "Wifipool · Chatbot", size=11, bold=True, color=INK_SOFT,
             anchor=MSO_ANCHOR.MIDDLE)

    if kicker:
        add_text(slide, Inches(0.85), Inches(0.95), Inches(10), Inches(0.32),
                 kicker, size=11, bold=True, color=PRIMARY,
                 anchor=MSO_ANCHOR.TOP)
    if title:
        add_text(slide, Inches(0.85), Inches(1.20), Inches(11.6), Inches(0.95),
                 title, size=34, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.05)

    # Footer
    if page_num is not None and total is not None:
        add_text(slide, Inches(0.45), Inches(7.05), Inches(6), Inches(0.3),
                 "Wifipool — Achter de schermen", size=9, color=INK_SOFT)
        add_text(slide, Inches(7), Inches(7.05), Inches(6), Inches(0.3),
                 f"{page_num} / {total}", size=9, color=INK_SOFT,
                 align=PP_ALIGN.RIGHT)


# ────────────────────────────────────────────────────────────────────
# Slides
# ────────────────────────────────────────────────────────────────────
TOTAL = 18
slides_data_built = []

def cover():
    s = prs.slides.add_slide(BLANK)
    # Background gradient-ish: navy with two soft circles
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    add_oval(s, Inches(-2), Inches(-2), Inches(7), Inches(7), OCEAN_700)
    add_oval(s, Inches(8.5), Inches(-1.5), Inches(6.5), Inches(6.5), OCEAN_500)
    add_oval(s, Inches(9), Inches(4.5), Inches(4), Inches(4), PRIMARY)

    # Logo block
    add_oval(s, Inches(0.85), Inches(0.85), Inches(0.7), Inches(0.7), WHITE)
    add_text(s, Inches(0.85), Inches(0.85), Inches(0.7), Inches(0.7),
             "🌊", size=24, color=OCEAN_500,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")
    add_text(s, Inches(1.7), Inches(0.85), Inches(5), Inches(0.7),
             "Wifipool · Chatbot", size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    # Eyebrow
    add_text(s, Inches(0.85), Inches(2.6), Inches(11), Inches(0.5),
             "ACHTER DE SCHERMEN", size=14, bold=True, color=OCEAN_400)
    # Title (huge)
    add_text(s, Inches(0.85), Inches(3.05), Inches(11), Inches(2.0),
             "Hoe onze chatbot\nleeft, leert en groeit",
             size=58, bold=True, color=WHITE, line_spacing=1.05)
    # Sub
    add_text(s, Inches(0.85), Inches(5.45), Inches(11), Inches(0.9),
             "Een visuele rondleiding doorheen het beheer, de werking en\nde updates van de Wifipool kennisassistent.",
             size=18, color=OCEAN_100, line_spacing=1.3)
    # Tag
    add_round(s, Inches(0.85), Inches(6.55), Inches(2.6), Inches(0.45),
              OCEAN_500, radius=0.5)
    add_text(s, Inches(0.85), Inches(6.55), Inches(2.6), Inches(0.45),
             "🚀  versie 2.1 · 2026",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")


def agenda():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="INHOUD", title="Wat ga je ontdekken?",
                page_num=2, total=TOTAL)

    items = [
        ("01", "🛠️", "Hoe wijzig je de FAQ?",
         "Het admin-dashboard, snel toevoegen, Excel bulk-edit, tabs per categorie."),
        ("02", "🧠", "Hoe werkt de chatbot?",
         "De architectuur, slim zoeken, meertalige antwoorden, synoniemen."),
        ("03", "🚀", "Hoe rollen we updates uit?",
         "Live wijzigingen, backups, monitoring en veiligheidsnet."),
    ]
    top = Inches(2.4)
    h = Inches(1.45)
    gap = Inches(0.18)
    for i, (num, emoji, title, sub) in enumerate(items):
        y = top + (h + gap) * i
        add_round(s, Inches(0.85), y, Inches(11.6), h, WHITE, line=LINE, radius=0.18)
        # Number plate
        add_round(s, Inches(1.1), y + Inches(0.22), Inches(1.0), Inches(1.0),
                  PRIMARY, radius=0.25)
        add_text(s, Inches(1.1), y + Inches(0.22), Inches(1.0), Inches(1.0),
                 num, size=26, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Emoji
        add_text(s, Inches(2.4), y + Inches(0.30), Inches(0.8), Inches(0.85),
                 emoji, size=34, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Segoe UI Emoji")
        # Title
        add_text(s, Inches(3.3), y + Inches(0.26), Inches(8.8), Inches(0.5),
                 title, size=20, bold=True, color=NAVY)
        add_text(s, Inches(3.3), y + Inches(0.78), Inches(8.8), Inches(0.6),
                 sub, size=13, color=INK_SOFT, line_spacing=1.3)


def divider(num, kicker, big_title, sub, accent=PRIMARY, page_num=0):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    # Decorative
    add_oval(s, Inches(-3), Inches(2), Inches(8), Inches(8), OCEAN_700)
    add_oval(s, Inches(9), Inches(-2), Inches(6), Inches(6), accent)
    add_oval(s, Inches(10.5), Inches(4.5), Inches(3.5), Inches(3.5), OCEAN_500)

    add_text(s, Inches(0.85), Inches(2.4), Inches(11), Inches(0.5),
             f"DEEL {num}", size=14, bold=True, color=OCEAN_400)
    add_text(s, Inches(0.85), Inches(2.85), Inches(11), Inches(0.5),
             kicker, size=14, bold=True, color=WHITE)
    add_text(s, Inches(0.85), Inches(3.45), Inches(11), Inches(2.4),
             big_title, size=64, bold=True, color=WHITE, line_spacing=1.0)
    add_text(s, Inches(0.85), Inches(5.85), Inches(11), Inches(0.7),
             sub, size=16, color=OCEAN_100, line_spacing=1.3)
    # Footer
    add_text(s, Inches(7), Inches(7.05), Inches(6), Inches(0.3),
             f"{page_num} / {TOTAL}", size=9, color=OCEAN_100,
             align=PP_ALIGN.RIGHT)


def slide_dashboard_overview():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 1 · DASHBOARD",
                title="Eén plek voor al je vragen",
                page_num=4, total=TOTAL)

    # Left panel
    add_round(s, Inches(0.85), Inches(2.3), Inches(7.5), Inches(4.5),
              SUBTLE, line=LINE, radius=0.06)
    add_text(s, Inches(1.1), Inches(2.55), Inches(7), Inches(0.5),
             "Het admin-dashboard", size=20, bold=True, color=NAVY)
    add_text(s, Inches(1.1), Inches(3.05), Inches(7), Inches(0.5),
             "Veilig achter een wachtwoord — voor jou en je team.",
             size=13, color=INK_SOFT)

    feats = [
        ("📊", "Overview",      "De belangrijkste cijfers in één oogopslag."),
        ("📚", "FAQ per categorie", "Bewerk vragen per thema — met tabs."),
        ("❓", "Onbeantwoord",  "Wat de chatbot niet kon beantwoorden."),
        ("📥", "Excel-import",  "Bulk-wijzigingen via je vertrouwde Excel."),
    ]
    for i, (e, t, d) in enumerate(feats):
        y = Inches(3.7) + Inches(0.7) * i
        add_text(s, Inches(1.2), y, Inches(0.4), Inches(0.5), e,
                 size=18, font="Segoe UI Emoji",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.7), y, Inches(2.0), Inches(0.5),
                 t, size=14, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.7), y, Inches(4.5), Inches(0.5),
                 d, size=12, color=INK_SOFT,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Right "screenshot" mock
    add_round(s, Inches(8.7), Inches(2.3), Inches(3.95), Inches(4.5),
              NAVY, radius=0.06)
    # title bar
    add_round(s, Inches(8.85), Inches(2.45), Inches(3.65), Inches(0.4),
              OCEAN_700, radius=0.25)
    add_oval(s, Inches(8.95), Inches(2.55), Inches(0.18), Inches(0.18), RED)
    add_oval(s, Inches(9.18), Inches(2.55), Inches(0.18), Inches(0.18), ACCENT)
    add_oval(s, Inches(9.41), Inches(2.55), Inches(0.18), Inches(0.18), GREEN)
    # Sidebar
    add_round(s, Inches(8.85), Inches(2.95), Inches(1.05), Inches(3.75),
              OCEAN_700, radius=0.05)
    items = ["📊 Overview", "📚 FAQ", "❓ Gaps"]
    for i, it in enumerate(items):
        add_round(s, Inches(8.95), Inches(3.10) + Inches(0.45) * i,
                  Inches(0.85), Inches(0.35),
                  OCEAN_500 if i == 1 else OCEAN_700, radius=0.4)
        add_text(s, Inches(8.95), Inches(3.10) + Inches(0.45) * i,
                 Inches(0.85), Inches(0.35),
                 it, size=8, bold=(i == 1), color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Segoe UI Emoji")
    # Main mock
    add_round(s, Inches(10.0), Inches(2.95), Inches(2.5), Inches(0.55),
              WHITE, radius=0.15)
    add_text(s, Inches(10.0), Inches(2.95), Inches(2.5), Inches(0.55),
             "FAQ per categorie", size=10, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cats = [("Algemeen", True), ("Wifipool", False), ("Display", False)]
    for i, (c, active) in enumerate(cats):
        add_round(s, Inches(10.0) + Inches(0.85) * i, Inches(3.6),
                  Inches(0.78), Inches(0.32),
                  PRIMARY if active else WHITE,
                  line=None if active else LINE, radius=0.5)
        add_text(s, Inches(10.0) + Inches(0.85) * i, Inches(3.6),
                 Inches(0.78), Inches(0.32),
                 c, size=8, bold=True,
                 color=WHITE if active else NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Rows
    for i in range(4):
        add_round(s, Inches(10.0), Inches(4.05) + Inches(0.55) * i,
                  Inches(2.45), Inches(0.45), WHITE, radius=0.15)
        add_round(s, Inches(10.15), Inches(4.15) + Inches(0.55) * i,
                  Inches(0.6), Inches(0.25), OCEAN_100, radius=0.5)
        add_round(s, Inches(10.85), Inches(4.13) + Inches(0.55) * i,
                  Inches(1.5), Inches(0.1), LINE, radius=0.5)
        add_round(s, Inches(10.85), Inches(4.30) + Inches(0.55) * i,
                  Inches(1.0), Inches(0.08), LINE, radius=0.5)


def slide_quick_add():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 1 · OPTIE 1",
                title="Snel een nieuwe vraag toevoegen",
                page_num=5, total=TOTAL)

    add_text(s, Inches(0.85), Inches(2.20), Inches(11.6), Inches(0.5),
             "Voor 1 tot 5 vragen — direct geleerd, in alle talen.",
             size=14, color=INK_SOFT)

    steps = [
        ("1", "Klik op 'Nieuwe vraag'", "Een eenvoudig formulier opent."),
        ("2", "Typ de vraag en het antwoord",
         "Voeg eventueel een YouTube-link of foto toe."),
        ("3", "Klik op opslaan", "De chatbot weet het meteen — voor elke taal."),
    ]
    box_w = Inches(3.85)
    box_h = Inches(3.2)
    gap = Inches(0.18)
    x0 = Inches(0.85)
    y0 = Inches(3.0)
    for i, (num, t, d) in enumerate(steps):
        x = x0 + (box_w + gap) * i
        add_round(s, x, y0, box_w, box_h, WHITE, line=LINE, radius=0.06)
        # accent strip
        add_round(s, x, y0, box_w, Inches(0.08), PRIMARY, radius=0)
        # number plate
        add_round(s, x + Inches(0.3), y0 + Inches(0.4),
                  Inches(0.7), Inches(0.7), OCEAN_50, radius=0.5)
        add_text(s, x + Inches(0.3), y0 + Inches(0.4),
                 Inches(0.7), Inches(0.7),
                 num, size=22, bold=True, color=PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.3), y0 + Inches(1.3),
                 box_w - Inches(0.6), Inches(0.7),
                 t, size=16, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y0 + Inches(2.0),
                 box_w - Inches(0.6), Inches(1.0),
                 d, size=12, color=INK_SOFT, line_spacing=1.4)

    add_round(s, Inches(0.85), Inches(6.40), Inches(11.6), Inches(0.55),
              OCEAN_50, line=OCEAN_400, radius=0.3)
    add_text(s, Inches(0.85), Inches(6.40), Inches(11.6), Inches(0.55),
             "💡  Tip: foto's en links kunnen later via Excel verfijnd worden.",
             size=12, bold=True, color=OCEAN_700,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")


def slide_excel_bulk():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 1 · OPTIE 2",
                title="Bulk-wijzigingen via Excel",
                page_num=6, total=TOTAL)

    add_text(s, Inches(0.85), Inches(2.20), Inches(11.6), Inches(0.5),
             "Voor grote bewerkingen — gebruik je vertrouwde Excel-sheet.",
             size=14, color=INK_SOFT)

    # Three big steps with arrows
    items = [
        ("📥", "Download", "De volledige FAQ als één Excel-bestand."),
        ("✍️", "Bewerk",   "Pas vragen, antwoorden en vertalingen aan."),
        ("📤", "Upload",   "De chatbot leert meteen van het nieuwe bestand."),
    ]
    box_w = Inches(3.4); box_h = Inches(3.0)
    gap_x = Inches(0.6)
    x0 = Inches(0.95); y0 = Inches(3.05)
    for i, (e, t, d) in enumerate(items):
        x = x0 + (box_w + gap_x) * i
        add_round(s, x, y0, box_w, box_h, NAVY, radius=0.06)
        # icon
        add_round(s, x + Inches(0.3), y0 + Inches(0.3),
                  Inches(0.9), Inches(0.9), OCEAN_700, radius=0.25)
        add_text(s, x + Inches(0.3), y0 + Inches(0.3),
                 Inches(0.9), Inches(0.9),
                 e, size=34, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Segoe UI Emoji")
        add_text(s, x + Inches(0.3), y0 + Inches(1.4),
                 box_w - Inches(0.6), Inches(0.6),
                 t, size=22, bold=True, color=WHITE)
        add_text(s, x + Inches(0.3), y0 + Inches(2.05),
                 box_w - Inches(0.6), Inches(0.8),
                 d, size=12, color=OCEAN_100, line_spacing=1.4)

        if i < 2:
            ax = x + box_w + Inches(0.05)
            ay = y0 + box_h / 2 - Inches(0.2)
            add_text(s, ax, ay, Inches(0.5), Inches(0.4), "➜",
                     size=28, bold=True, color=PRIMARY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_round(s, Inches(0.85), Inches(6.40), Inches(11.6), Inches(0.55),
              SUBTLE, line=LINE, radius=0.3)
    add_text(s, Inches(0.85), Inches(6.40), Inches(11.6), Inches(0.55),
             "🛟  Veilig: van elk geüpload bestand bewaren we automatisch een backup.",
             size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")


def slide_categories_tabs():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 1 · NIEUW",
                title="FAQ per categorie — met tabs",
                page_num=7, total=TOTAL)

    add_text(s, Inches(0.85), Inches(2.20), Inches(11.6), Inches(0.5),
             "Springt direct naar het thema dat je wil bewerken.",
             size=14, color=INK_SOFT)

    # Mock tabs row
    tabs = [("All", 336, True), ("Algemeen", 42, False),
            ("Wifipool", 88, False), ("Display", 24, False),
            ("EPDM", 11, False), ("pH meting", 18, False),
            ("Zoutelectrolyse", 35, False)]
    x = Inches(0.85); y = Inches(2.95)
    add_round(s, x, y, Inches(11.6), Inches(0.65), WHITE, line=LINE, radius=0.18)
    cur_x = x + Inches(0.18); cur_y = y + Inches(0.13)
    for label, n, active in tabs:
        w = Inches(0.18) + Inches(0.12) * len(label) + Inches(0.55)
        add_round(s, cur_x, cur_y, w, Inches(0.4),
                  PRIMARY if active else SUBTLE,
                  line=None if active else LINE, radius=0.5)
        add_text(s, cur_x, cur_y, w - Inches(0.35), Inches(0.4),
                 label, size=11, bold=True,
                 color=WHITE if active else NAVY,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # count pill
        add_round(s, cur_x + w - Inches(0.32), cur_y + Inches(0.08),
                  Inches(0.28), Inches(0.24),
                  WHITE if active else WHITE, radius=0.5)
        add_text(s, cur_x + w - Inches(0.32), cur_y + Inches(0.08),
                 Inches(0.28), Inches(0.24),
                 str(n), size=8, bold=True,
                 color=PRIMARY if active else INK_SOFT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += w + Inches(0.07)

    # Mock table
    add_round(s, Inches(0.85), Inches(3.85), Inches(11.6), Inches(2.65),
              WHITE, line=LINE, radius=0.06)
    headers = ["Categorie", "Vraag", "Antwoord", ""]
    widths = [Inches(2.0), Inches(4.5), Inches(4.0), Inches(1.1)]
    hx = Inches(0.85)
    add_round(s, Inches(0.85), Inches(3.85), Inches(11.6), Inches(0.4),
              SUBTLE, radius=0.06)
    for h, w in zip(headers, widths):
        add_text(s, hx + Inches(0.18), Inches(3.85), w, Inches(0.4),
                 h, size=10, bold=True, color=INK_SOFT,
                 anchor=MSO_ANCHOR.MIDDLE)
        hx += w
    # Rows
    samples = [
        ("Wifipool", "Hoe reset ik Gen 2?", "Druk 5s op de knop tot de led knippert…"),
        ("Display", "Display blijft zwart na opstart",
         "Controleer de voeding en de modulekabel…"),
        ("Algemeen", "Waar vind ik de handleidingen?",
         "Op beniferro.eu en op ons YouTube-kanaal."),
    ]
    for i, (c, q, a) in enumerate(samples):
        ry = Inches(4.30) + Inches(0.55) * i
        rx = Inches(0.85)
        add_round(s, rx + Inches(0.18), ry + Inches(0.07),
                  Inches(1.4), Inches(0.32),
                  OCEAN_100, radius=0.5)
        add_text(s, rx + Inches(0.18), ry + Inches(0.07),
                 Inches(1.4), Inches(0.32),
                 c, size=9, bold=True, color=OCEAN_700,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(2.18), ry + Inches(0.05),
                 Inches(4.3), Inches(0.4),
                 q, size=11, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(6.7), ry + Inches(0.05),
                 Inches(3.8), Inches(0.4),
                 a, size=11, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE)

    add_round(s, Inches(0.85), Inches(6.65), Inches(11.6), Inches(0.45),
              OCEAN_50, line=OCEAN_400, radius=0.3)
    add_text(s, Inches(0.85), Inches(6.65), Inches(11.6), Inches(0.45),
             "✨  Nieuw: download óók een Excel met één tabblad per categorie.",
             size=12, bold=True, color=OCEAN_700,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")


# ── DEEL 2 slides ─────────────────────────────────────────────────────
def slide_arch_overview():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 2 · ARCHITECTUUR",
                title="Wat gebeurt er als iemand een vraag stelt?",
                page_num=9, total=TOTAL)

    # Pipeline
    steps = [
        ("👤", "Klant",        "Stelt een vraag\nop de website"),
        ("🔍", "Begrijpen",    "Detecteert taal\nen bedoeling"),
        ("📚", "Zoeken",       "Doorzoekt de\nFAQ slim"),
        ("🤖", "Antwoorden",   "Schrijft het juiste\nantwoord op maat"),
        ("💬", "Tonen",        "Met foto's, video\nen vertaling"),
    ]
    n = len(steps)
    box_w = Inches(2.15); box_h = Inches(2.55)
    total_w = box_w * n + Inches(0.35) * (n - 1)
    x0 = (SLIDE_W - total_w) / 2
    y0 = Inches(2.85)
    for i, (e, t, d) in enumerate(steps):
        x = x0 + (box_w + Inches(0.35)) * i
        # arrow before
        if i > 0:
            ax = x - Inches(0.42)
            add_text(s, ax, y0 + box_h/2 - Inches(0.3),
                     Inches(0.5), Inches(0.5), "➜",
                     size=28, bold=True, color=PRIMARY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round(s, x, y0, box_w, box_h, WHITE, line=LINE, radius=0.10)
        add_round(s, x, y0, box_w, Inches(0.08), PRIMARY, radius=0)
        # circle icon
        add_oval(s, x + box_w/2 - Inches(0.55), y0 + Inches(0.35),
                 Inches(1.1), Inches(1.1), OCEAN_50)
        add_text(s, x + box_w/2 - Inches(0.55), y0 + Inches(0.35),
                 Inches(1.1), Inches(1.1),
                 e, size=34, font="Segoe UI Emoji",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # title
        add_text(s, x, y0 + Inches(1.55),
                 box_w, Inches(0.4),
                 t, size=15, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        # desc
        add_text(s, x + Inches(0.18), y0 + Inches(1.95),
                 box_w - Inches(0.36), Inches(0.6),
                 d, size=11, color=INK_SOFT,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

    add_round(s, Inches(0.85), Inches(5.95), Inches(11.6), Inches(0.85),
              SUBTLE, line=LINE, radius=0.3)
    add_text(s, Inches(1.1), Inches(5.95), Inches(11.1), Inches(0.85),
             "⚡  Het hele proces gebeurt in milliseconden — alles draait in de cloud,\nzonder dat de bezoeker iets merkt van de complexiteit erachter.",
             size=13, color=NAVY, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji", line_spacing=1.4)


def slide_smart_search():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 2 · SLIM ZOEKEN",
                title="Zoeken zoals een mens",
                page_num=10, total=TOTAL)

    # Left: question
    add_round(s, Inches(0.85), Inches(2.45), Inches(5.8), Inches(2.0),
              OCEAN_50, line=OCEAN_400, radius=0.06)
    add_text(s, Inches(1.15), Inches(2.65), Inches(5.4), Inches(0.4),
             "DE VRAAG VAN DE KLANT", size=11, bold=True, color=OCEAN_700)
    add_text(s, Inches(1.15), Inches(3.05), Inches(5.4), Inches(1.2),
             "« Mijn doseerpomp blijft\ndruppelen, wat doe ik fout? »",
             size=22, bold=True, color=NAVY, line_spacing=1.2)

    # Right: thinking bubble
    add_round(s, Inches(7.0), Inches(2.45), Inches(5.5), Inches(4.4),
              NAVY, radius=0.06)
    add_text(s, Inches(7.3), Inches(2.65), Inches(5.0), Inches(0.4),
             "DE CHATBOT DENKT…", size=11, bold=True, color=OCEAN_400)
    bullets = [
        ("✓", "Begrijpt 'doseerpomp' = synoniem voor injector"),
        ("✓", "Herkent de bedoeling: probleemoplossing"),
        ("✓", "Vergelijkt met honderden FAQ's tegelijk"),
        ("✓", "Vindt de meest waarschijnlijke match"),
        ("✓", "Controleert dat het antwoord echt past"),
    ]
    for i, (m, t) in enumerate(bullets):
        y = Inches(3.15) + Inches(0.62) * i
        add_oval(s, Inches(7.3), y + Inches(0.05),
                 Inches(0.32), Inches(0.32), GREEN)
        add_text(s, Inches(7.3), y + Inches(0.05),
                 Inches(0.32), Inches(0.32), m,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(7.78), y, Inches(4.7), Inches(0.5),
                 t, size=13, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Bottom: result
    add_round(s, Inches(0.85), Inches(5.0), Inches(5.8), Inches(1.85),
              WHITE, line=LINE, radius=0.06)
    add_text(s, Inches(1.15), Inches(5.15), Inches(5.4), Inches(0.4),
             "✓ ANTWOORD GEVONDEN", size=11, bold=True, color=GREEN,
             font="Segoe UI Emoji")
    add_text(s, Inches(1.15), Inches(5.55), Inches(5.4), Inches(1.2),
             "Controleer of de terugslagklep van de\npomp niet vuil is en eventueel de\nO-ringen vervangen.",
             size=13, color=INK_SOFT, line_spacing=1.35)


def slide_multilang():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 2 · MEERTALIG",
                title="Eén vraag, vier talen",
                page_num=11, total=TOTAL)

    add_text(s, Inches(0.85), Inches(2.20), Inches(11.6), Inches(0.5),
             "De chatbot herkent automatisch de taal van de bezoeker en antwoordt navenant.",
             size=14, color=INK_SOFT)

    flags = [
        ("🇳🇱", "Nederlands", "Brontaal — alle vragen worden eerst hier beheerd."),
        ("🇬🇧", "English",    "Antwoorden vooraf vertaald, kant-en-klaar."),
        ("🇫🇷", "Français",   "Idem — even snel als in het Nederlands."),
        ("🇩🇪", "Deutsch",    "Zelfde kwaliteit, elk dialect ondersteund."),
    ]
    box_w = Inches(2.85); box_h = Inches(3.5)
    gap = Inches(0.13); x0 = Inches(0.85); y0 = Inches(3.0)
    for i, (f, lang, d) in enumerate(flags):
        x = x0 + (box_w + gap) * i
        add_round(s, x, y0, box_w, box_h, WHITE, line=LINE, radius=0.08)
        # flag big
        add_text(s, x, y0 + Inches(0.4),
                 box_w, Inches(1.2),
                 f, size=64, font="Segoe UI Emoji",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # name
        add_text(s, x, y0 + Inches(1.7),
                 box_w, Inches(0.5),
                 lang, size=20, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        # desc
        add_text(s, x + Inches(0.3), y0 + Inches(2.25),
                 box_w - Inches(0.6), Inches(1.1),
                 d, size=12, color=INK_SOFT,
                 align=PP_ALIGN.CENTER, line_spacing=1.4)

    add_round(s, Inches(0.85), Inches(6.65), Inches(11.6), Inches(0.45),
              OCEAN_50, line=OCEAN_400, radius=0.3)
    add_text(s, Inches(0.85), Inches(6.65), Inches(11.6), Inches(0.45),
             "🌐  Voeg in één taal toe → de andere drie volgen automatisch.",
             size=12, bold=True, color=OCEAN_700,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")


def slide_synonyms():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 2 · SYNONIEMEN",
                title="Verschillende woorden, zelfde antwoord",
                page_num=12, total=TOTAL)

    add_text(s, Inches(0.85), Inches(2.20), Inches(11.6), Inches(0.5),
             "Klanten gebruiken hun eigen woordenschat. De chatbot ook.",
             size=14, color=INK_SOFT)

    # Visual: bunch of bubbles converging
    bubbles_left = [
        ("doseerpomp", Inches(1.0), Inches(3.1)),
        ("injector",   Inches(2.6), Inches(3.5)),
        ("doseerunit", Inches(1.4), Inches(4.3)),
        ("pomp",       Inches(2.8), Inches(4.7)),
        ("dosing pump", Inches(1.0), Inches(5.5)),
    ]
    for label, x, y in bubbles_left:
        w = Inches(0.18 + 0.12 * len(label) + 0.4)
        add_round(s, x, y, w, Inches(0.55),
                  WHITE, line=OCEAN_400, radius=0.5)
        add_text(s, x, y, w, Inches(0.55),
                 label, size=12, bold=True, color=OCEAN_700,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Arrow
    add_text(s, Inches(5.5), Inches(4.05), Inches(2.0), Inches(0.6),
             "➜", size=44, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Right: single concept
    add_round(s, Inches(8.0), Inches(3.4), Inches(4.5), Inches(2.5),
              NAVY, radius=0.08)
    add_text(s, Inches(8.0), Inches(3.55), Inches(4.5), Inches(0.5),
             "🎯  ÉÉN BEGRIP", size=12, bold=True, color=OCEAN_400,
             align=PP_ALIGN.CENTER, font="Segoe UI Emoji")
    add_text(s, Inches(8.0), Inches(4.05), Inches(4.5), Inches(0.7),
             "Doseersysteem", size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.3), Inches(4.85), Inches(3.9), Inches(1.0),
             "Eén juist antwoord —\nongeacht hoe de klant het noemt.",
             size=13, color=OCEAN_100,
             align=PP_ALIGN.CENTER, line_spacing=1.4)

    add_round(s, Inches(0.85), Inches(6.55), Inches(11.6), Inches(0.45),
              SUBTLE, line=LINE, radius=0.3)
    add_text(s, Inches(0.85), Inches(6.55), Inches(11.6), Inches(0.45),
             "✏️  Synoniemen kunnen we centraal beheren in een speciaal Excel-tabblad.",
             size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")


def slide_data_used():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 2 · KENNISBANK",
                title="Wat zit er allemaal in de chatbot?",
                page_num=13, total=TOTAL)

    stats = [
        ("336+", "Vragen & antwoorden", OCEAN_500),
        ("4",     "Talen ondersteund",   PRIMARY),
        ("17",    "Categorieën",         OCEAN_700),
        ("∞",     "Synoniem-varianten",  ACCENT),
    ]
    box_w = Inches(2.85); box_h = Inches(2.4)
    gap = Inches(0.13); x0 = Inches(0.85); y0 = Inches(2.7)
    for i, (n, l, c) in enumerate(stats):
        x = x0 + (box_w + gap) * i
        add_round(s, x, y0, box_w, box_h, WHITE, line=LINE, radius=0.08)
        add_round(s, x, y0, box_w, Inches(0.08), c, radius=0)
        add_text(s, x, y0 + Inches(0.6),
                 box_w, Inches(1.2),
                 n, size=64, bold=True, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y0 + Inches(1.75),
                 box_w, Inches(0.5),
                 l, size=14, bold=True, color=INK_SOFT,
                 align=PP_ALIGN.CENTER)

    # Extra row
    extras = [
        ("📷", "Foto's per vraag",        "Visuele uitleg waar nodig."),
        ("🎬", "Video-links",              "Naar het YouTube-kanaal van Beniferro."),
        ("🔄", "Alternatieve formuleringen", "Voor als de klant het anders zegt."),
    ]
    bx = Inches(0.85); by = Inches(5.4)
    bw = Inches(3.85); bh = Inches(1.5)
    for i, (e, t, d) in enumerate(extras):
        x = bx + (bw + Inches(0.18)) * i
        add_round(s, x, by, bw, bh, SUBTLE, radius=0.06)
        add_text(s, x + Inches(0.25), by + Inches(0.28),
                 Inches(0.5), Inches(0.5),
                 e, size=24, font="Segoe UI Emoji",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), by + Inches(0.22),
                 bw - Inches(1.0), Inches(0.45),
                 t, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.85), by + Inches(0.7),
                 bw - Inches(1.0), Inches(0.7),
                 d, size=11, color=INK_SOFT, line_spacing=1.3)


# ── DEEL 3 slides ─────────────────────────────────────────────────────
def slide_continuous_updates():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 3 · UPDATES",
                title="Nieuwe vragen, meteen live",
                page_num=15, total=TOTAL)

    add_text(s, Inches(0.85), Inches(2.20), Inches(11.6), Inches(0.5),
             "Geen onderhoudsvenster, geen downtime, geen technische tussenkomst nodig.",
             size=14, color=INK_SOFT)

    # Timeline
    steps = [
        ("0s",   "Je drukt op opslaan", OCEAN_500),
        ("1s",   "Server herlaadt de FAQ", OCEAN_500),
        ("2s",   "Nieuwe vraag staat in de zoekindex", PRIMARY),
        ("3s",   "Klanten krijgen het nieuwe antwoord", GREEN),
    ]
    # Horizontal line
    add_rect(s, Inches(1.5), Inches(4.4), Inches(10.4), Inches(0.04), LINE)
    n = len(steps)
    for i, (label, t, c) in enumerate(steps):
        x = Inches(1.5) + (Inches(10.4) / (n - 1)) * i - Inches(0.5)
        add_oval(s, x, Inches(4.15), Inches(0.55), Inches(0.55), c)
        add_text(s, x, Inches(4.15), Inches(0.55), Inches(0.55), "✓",
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # label below
        add_text(s, x - Inches(1.0), Inches(4.85), Inches(2.5), Inches(0.4),
                 label, size=18, bold=True, color=c,
                 align=PP_ALIGN.CENTER)
        add_text(s, x - Inches(1.4), Inches(5.25), Inches(3.3), Inches(0.7),
                 t, size=12, color=INK_SOFT,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

    # Bottom callout
    add_round(s, Inches(0.85), Inches(6.30), Inches(11.6), Inches(0.65),
              NAVY, radius=0.3)
    add_text(s, Inches(0.85), Inches(6.30), Inches(11.6), Inches(0.65),
             "🚀  De klant aan de andere kant merkt nul onderbreking.",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font="Segoe UI Emoji")


def slide_backup_safety():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 3 · VEILIGHEIDSNET",
                title="Backups en herstel",
                page_num=16, total=TOTAL)

    items = [
        ("💾", "Automatische backups",
         "Bij elke Excel-upload bewaren we de vorige versie.\nDe laatste 3 versies blijven beschikbaar."),
        ("🔐", "Beveiligde toegang",
         "Het dashboard zit achter een wachtwoord.\nAlleen jij en je team kunnen wijzigen."),
        ("⏪", "Snel terugdraaien",
         "Foutje gemaakt? In één klik herstel je\nde vorige versie van de FAQ."),
        ("🛡️", "Foutbestendigheid",
         "Als iets misgaat tijdens een upload,\nblijft de oude versie automatisch actief."),
    ]
    bx = Inches(0.85); by = Inches(2.6)
    bw = Inches(5.7); bh = Inches(1.95); gx = Inches(0.2); gy = Inches(0.18)
    for i, (e, t, d) in enumerate(items):
        col = i % 2; row = i // 2
        x = bx + (bw + gx) * col
        y = by + (bh + gy) * row
        add_round(s, x, y, bw, bh, WHITE, line=LINE, radius=0.06)
        add_round(s, x + Inches(0.25), y + Inches(0.3),
                  Inches(1.0), Inches(1.0), OCEAN_50, radius=0.5)
        add_text(s, x + Inches(0.25), y + Inches(0.3),
                 Inches(1.0), Inches(1.0), e,
                 size=34, font="Segoe UI Emoji",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.45), y + Inches(0.3),
                 bw - Inches(1.65), Inches(0.5),
                 t, size=18, bold=True, color=NAVY)
        add_text(s, x + Inches(1.45), y + Inches(0.85),
                 bw - Inches(1.65), Inches(1.0),
                 d, size=12, color=INK_SOFT, line_spacing=1.4)


def slide_monitoring():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="DEEL 3 · MONITORING",
                title="We zien wat de klanten vragen",
                page_num=17, total=TOTAL)

    add_text(s, Inches(0.85), Inches(2.20), Inches(11.6), Inches(0.5),
             "Het dashboard toont continu wat goed loopt en wat verbetering nodig heeft.",
             size=14, color=INK_SOFT)

    # Big "metric" cards
    metrics = [
        ("🔥", "Top vragen",
         "De 10 meest gestelde vragen — handig om je\nFAQ te prioriteren."),
        ("❓", "Onbeantwoorde vragen",
         "Wat de chatbot niet kon helpen — kandidaten om\nde kennisbank uit te breiden."),
        ("🌐", "Talen-verdeling",
         "Welke taal je klanten gebruiken —\nfocus op de juiste vertalingen."),
    ]
    bx = Inches(0.85); by = Inches(2.95)
    bw = Inches(3.85); bh = Inches(3.5)
    for i, (e, t, d) in enumerate(metrics):
        x = bx + (bw + Inches(0.18)) * i
        add_round(s, x, by, bw, bh, WHITE, line=LINE, radius=0.08)
        add_round(s, x, by, bw, Inches(0.08), PRIMARY, radius=0)
        add_round(s, x + Inches(0.3), by + Inches(0.4),
                  Inches(1.0), Inches(1.0), PRIMARY_2, radius=0.5)
        add_text(s, x + Inches(0.3), by + Inches(0.4),
                 Inches(1.0), Inches(1.0),
                 e, size=34, font="Segoe UI Emoji",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.3), by + Inches(1.55),
                 bw - Inches(0.6), Inches(0.5),
                 t, size=18, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), by + Inches(2.05),
                 bw - Inches(0.6), Inches(1.3),
                 d, size=13, color=INK_SOFT, line_spacing=1.4)


def slide_recap():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, kicker="OVERZICHT",
                title="Wat heb je nu in handen?",
                page_num=18, total=TOTAL)

    items = [
        ("✅", "Een dashboard om alles te beheren — zonder ontwikkelaars."),
        ("✅", "Snel toevoegen of bulk-bewerken via Excel — jouw keuze."),
        ("✅", "Een meertalige chatbot die slim zoekt en netjes antwoordt."),
        ("✅", "Live updates en automatische backups."),
        ("✅", "Inzicht in wat klanten vragen — om continu te verbeteren."),
    ]
    by = Inches(2.6); bx = Inches(1.4); bw = Inches(10.5); bh = Inches(0.65)
    gap = Inches(0.18)
    for i, (e, t) in enumerate(items):
        y = by + (bh + gap) * i
        add_round(s, bx, y, bw, bh, SUBTLE, radius=0.3)
        add_oval(s, bx + Inches(0.15), y + Inches(0.13),
                 Inches(0.4), Inches(0.4), GREEN)
        add_text(s, bx + Inches(0.15), y + Inches(0.13),
                 Inches(0.4), Inches(0.4), "✓",
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, bx + Inches(0.7), y, bw - Inches(0.85), bh,
                 t, size=15, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    add_oval(s, Inches(-2), Inches(-2), Inches(8), Inches(8), OCEAN_700)
    add_oval(s, Inches(8), Inches(-2), Inches(7), Inches(7), PRIMARY)
    add_oval(s, Inches(9), Inches(4.5), Inches(4), Inches(4), OCEAN_500)

    add_text(s, Inches(0.85), Inches(2.5), Inches(11.6), Inches(1.0),
             "BEDANKT", size=18, bold=True, color=OCEAN_400)
    add_text(s, Inches(0.85), Inches(3.0), Inches(11.6), Inches(2.5),
             "Vragen?", size=128, bold=True, color=WHITE, line_spacing=1.0)
    add_text(s, Inches(0.85), Inches(5.7), Inches(11.6), Inches(0.7),
             "We helpen je graag verder met de volgende stappen.",
             size=18, color=OCEAN_100)


# ────────────────────────────────────────────────────────────────────
# Build deck
# ────────────────────────────────────────────────────────────────────
cover()                                                # 1
agenda()                                               # 2
divider(1, "DE FAQ BEHEREN", "Hoe wijzig je\nde kennisbank?",
        "Drie eenvoudige manieren om je chatbot te voeden.",
        accent=OCEAN_500, page_num=3)                  # 3
slide_dashboard_overview()                             # 4
slide_quick_add()                                      # 5
slide_excel_bulk()                                     # 6
slide_categories_tabs()                                # 7
divider(2, "ONDER DE MOTORKAP",
        "Hoe werkt de\nchatbot écht?",
        "Een blik op de logica zonder technisch jargon.",
        accent=PRIMARY, page_num=8)                    # 8
slide_arch_overview()                                  # 9
slide_smart_search()                                   # 10
slide_multilang()                                      # 11
slide_synonyms()                                       # 12
slide_data_used()                                      # 13
divider(3, "LEVEND SYSTEEM",
        "Updates &\nveiligheidsnet",
        "Hoe we evolueren zonder iets te breken.",
        accent=OCEAN_700, page_num=14)                 # 14
slide_continuous_updates()                             # 15
slide_backup_safety()                                  # 16
slide_monitoring()                                     # 17
slide_recap()                                          # 18 (recap as the official "what's there")
slide_thanks()                                         # 19 (closing — extra)

out = "Wifipool_Chatbot_Achter_de_Schermen.pptx"
prs.save(out)
print(f"[OK] Wrote {out} with {len(prs.slides)} slides.")
