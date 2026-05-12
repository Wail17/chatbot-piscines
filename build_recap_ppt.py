"""Genereer een korte 2-slide PowerPoint update over de chatbot voortgang.

In het Nederlands. Modern donker design met accent kleur Wifipool blauw.
Output: WifipoolChatbot_Update.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Design system ------------------------------------------------------------
COLOR_BG          = RGBColor(0x0B, 0x14, 0x22)   # diep navy
COLOR_BG_2        = RGBColor(0x12, 0x1E, 0x33)   # iets lichter navy voor cards
COLOR_ACCENT      = RGBColor(0x38, 0xBD, 0xF8)   # Wifipool cyaan/blauw
COLOR_ACCENT_DARK = RGBColor(0x0E, 0xA5, 0xE9)
COLOR_TEXT        = RGBColor(0xE5, 0xE7, 0xEB)   # bijna wit
COLOR_TEXT_DIM    = RGBColor(0x94, 0xA3, 0xB8)   # grijs
COLOR_OK          = RGBColor(0x4A, 0xDE, 0x80)   # groen check
COLOR_WARN        = RGBColor(0xFB, 0xBF, 0x24)   # geel volgende stap

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    set_solid(shp, color)
    if rounded:
        # Subtle corner radius
        shp.adjustments[0] = 0.06
    return shp


def add_text(slide, x, y, w, h, text, size=14, color=COLOR_TEXT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_accent_bar(slide, x, y, w=Inches(0.08), h=Inches(0.45)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_solid(bar, COLOR_ACCENT)


def add_chip(slide, x, y, w, label, color=COLOR_ACCENT):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34))
    chip.adjustments[0] = 0.5
    set_solid(chip, color)
    tf = chip.text_frame
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Consolas"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = COLOR_BG


def add_bullet(slide, x, y, w, h, icon, title, body):
    """Card-ish bullet: small accent square + bold title + dim body."""
    # icon block
    ic = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.42), Inches(0.42))
    ic.adjustments[0] = 0.25
    set_solid(ic, COLOR_BG_2)
    itf = ic.text_frame
    itf.margin_left = Emu(0); itf.margin_right = Emu(0)
    itf.margin_top = Emu(0); itf.margin_bottom = Emu(0)
    itf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ip = itf.paragraphs[0]
    ip.alignment = PP_ALIGN.CENTER
    irun = ip.add_run()
    irun.text = icon
    irun.font.size = Pt(16)
    irun.font.color.rgb = COLOR_ACCENT

    # text
    tx_x = x + Inches(0.6)
    tx_w = w - Inches(0.6)
    title_tb = add_text(slide, tx_x, y - Inches(0.02), tx_w, Inches(0.35),
                        title, size=14, color=COLOR_TEXT, bold=True)
    body_tb = add_text(slide, tx_x, y + Inches(0.33), tx_w, h - Inches(0.35),
                       body, size=11, color=COLOR_TEXT_DIM)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_solid(bg, COLOR_BG)
    return slide


def add_header(slide, eyebrow, title):
    add_text(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.32),
             eyebrow.upper(), size=11, color=COLOR_ACCENT, bold=True)
    add_text(slide, Inches(0.6), Inches(0.78), Inches(12), Inches(0.7),
             title, size=30, color=COLOR_TEXT, bold=True)
    # accent underline
    add_rect(slide, Inches(0.6), Inches(1.45), Inches(0.6), Emu(40000), COLOR_ACCENT)


def add_footer(slide, page_num, total=2):
    add_text(slide, Inches(0.6), Inches(7.0), Inches(6), Inches(0.3),
             "Wifipool / Beniferro · Chatbot AI · Wail Khamlichi",
             size=9, color=COLOR_TEXT_DIM)
    add_text(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
             f"{page_num} / {total}", size=9, color=COLOR_TEXT_DIM,
             align=PP_ALIGN.RIGHT)


# =============================================================================
# Slide 1 — Wat is er deze week opgelost
# =============================================================================

def build_slide_1(prs):
    s = new_slide(prs)
    add_header(s, "Voortgang  ·  Week 19",
               "Chatbot is opnieuw operationeel ✅")

    # subtitle
    add_text(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.4),
             "Wat we deze week gefixt en gebouwd hebben",
             size=15, color=COLOR_TEXT_DIM)

    # Tech chips
    chip_y = Inches(2.2)
    chips = [
        ("ANTHROPIC CLAUDE HAIKU 4.5", Inches(2.6)),
        ("PROMPT CACHING", Inches(1.65)),
        ("JSONL STORE", Inches(1.3)),
        ("FALLBACK CHAIN", Inches(1.55)),
        ("FASTAPI + UVICORN", Inches(1.8)),
    ]
    cx = Inches(0.6)
    for label, w in chips:
        add_chip(s, cx, chip_y, w, label)
        cx += w + Inches(0.15)

    # Cards row 1
    row1_y = Inches(2.95)
    col_w = Inches(6.05)
    col_h = Inches(1.95)
    # card 1 — API
    card1 = add_rect(s, Inches(0.6), row1_y, col_w, col_h, COLOR_BG_2, rounded=True)
    add_accent_bar(s, Inches(0.6), row1_y + Inches(0.25))
    add_text(s, Inches(0.85), row1_y + Inches(0.2), col_w - Inches(0.4), Inches(0.4),
             "🔧  API gefixt", size=18, color=COLOR_TEXT, bold=True)
    add_text(s, Inches(0.85), row1_y + Inches(0.7), col_w - Inches(0.4), col_h - Inches(0.8),
             "De Anthropic Claude API kwam niet meer door (sleutel uitgeput op Render). "
             "Volledig hersteld met expliciete error-logging, een eerlijke fallback "
             "in de juiste taal, en prompt-caching van de volledige FAQ-context "
             "(335 entries) voor lage latency en lage kost per vraag.",
             size=12, color=COLOR_TEXT_DIM)

    # card 2 — Dashboard CRUD
    card2_x = Inches(0.6) + col_w + Inches(0.25)
    card2 = add_rect(s, card2_x, row1_y, col_w, col_h, COLOR_BG_2, rounded=True)
    add_accent_bar(s, card2_x, row1_y + Inches(0.25))
    add_text(s, card2_x + Inches(0.25), row1_y + Inches(0.2), col_w - Inches(0.4), Inches(0.4),
             "🛠️  Dashboard met volledige CRUD-controle", size=18, color=COLOR_TEXT, bold=True)
    add_text(s, card2_x + Inches(0.25), row1_y + Inches(0.7), col_w - Inches(0.4), col_h - Inches(0.8),
             "De admin kan nu vragen / antwoorden / categorieën rechtstreeks beheren: "
             "create, read, update, delete via REST endpoints (/admin/faq). "
             "Wijzigingen worden in real-time naar Claude doorgegeven dankzij "
             "content-hash cache-invalidatie.",
             size=12, color=COLOR_TEXT_DIM)

    # Bottom row — bullets
    row2_y = Inches(5.05)
    bullets = [
        ("🧪", "Live FAQ Simulator",
         "Test elke vraag in elke taal, bekijk source, confidence en de gematchte FAQ-rij, "
         "spring meteen naar de editor."),
        ("📊", "Geverifieerde nauwkeurigheid",
         "97.3% correcte matches op 74 testvragen (NL/NL-paraphrase/EN/FR). "
         "Target 95% gehaald."),
        ("🔒", "Veilige fallback",
         "Bij API-fout krijgt de gebruiker een nette boodschap in zijn taal, geen onzin-antwoord."),
    ]
    bw = Inches(4.0)
    bx = Inches(0.6)
    for icon, t, b in bullets:
        add_bullet(s, bx, row2_y, bw, Inches(1.7), icon, t, b)
        bx += bw + Inches(0.2)

    add_footer(s, 1)
    return s


# =============================================================================
# Slide 2 — Wat volgt + status
# =============================================================================

def build_slide_2(prs):
    s = new_slide(prs)
    add_header(s, "Planning  ·  Week 20",
               "Volgende stap: live op beniferro.eu 🚀")

    add_text(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.4),
             "Wat we volgende week doen voor de finale oplevering",
             size=15, color=COLOR_TEXT_DIM)

    # Big "Next" cards
    card_y = Inches(2.3)
    col_w = Inches(4.0)
    col_h = Inches(2.7)
    cols = [
        ("01", "🌐  Live op beniferro.eu",
         "De chatbot wordt geïntegreerd op de productiewebsite van Beniferro. "
         "Embed-versies (chatbot-demo-optimized.html) zijn al klaar, "
         "alleen nog koppelen aan het juiste subdomein en SSL valideren.",
         COLOR_ACCENT),
        ("02", "📚  Documentatie verbeteren",
         "README, deploy-gids, onboarding voor de admin (hoe een FAQ-rij toevoegen, "
         "hoe de Simulator gebruiken, hoe de Excel-export draaien). "
         "Eén bron van waarheid voor het hele team.",
         COLOR_ACCENT),
        ("03", "🧹  Code opkuisen",
         "Refactor van enkele legacy paden (oude rag.py / rag_pure.py varianten), "
         "dode imports verwijderen, type-hints aanvullen en de testsuite "
         "uitbreiden tot 50+ live testen.",
         COLOR_ACCENT),
    ]
    cx = Inches(0.6)
    for num, title, body, color in cols:
        c = add_rect(s, cx, card_y, col_w, col_h, COLOR_BG_2, rounded=True)
        # big number
        add_text(s, cx + Inches(0.25), card_y + Inches(0.15), col_w - Inches(0.5), Inches(0.6),
                 num, size=28, color=color, bold=True, font="Consolas")
        # accent bar under number
        add_rect(s, cx + Inches(0.25), card_y + Inches(0.75), Inches(0.45), Emu(30000), color)
        # title
        add_text(s, cx + Inches(0.25), card_y + Inches(0.95), col_w - Inches(0.5), Inches(0.45),
                 title, size=15, color=COLOR_TEXT, bold=True)
        # body
        add_text(s, cx + Inches(0.25), card_y + Inches(1.45), col_w - Inches(0.5), col_h - Inches(1.55),
                 body, size=11, color=COLOR_TEXT_DIM)
        cx += col_w + Inches(0.2)

    # Status banner
    banner_y = Inches(5.5)
    banner = add_rect(s, Inches(0.6), banner_y, Inches(12.1), Inches(1.05), COLOR_BG_2, rounded=True)
    add_accent_bar(s, Inches(0.6), banner_y + Inches(0.18), h=Inches(0.7))
    add_text(s, Inches(0.85), banner_y + Inches(0.12), Inches(11.5), Inches(0.4),
             "📌  PROJECT STATUS",
             size=11, color=COLOR_ACCENT, bold=True)
    add_text(s, Inches(0.85), banner_y + Inches(0.4), Inches(11.5), Inches(0.6),
             "Project is bijna afgerond — backend stabiel, frontend embed werkt, "
             "97% accuracy, admin-tooling volledig. Resterend werk = deploy + docs + cleanup.",
             size=14, color=COLOR_TEXT, bold=True)

    add_footer(s, 2)
    return s


# =============================================================================
# Main
# =============================================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs)
    build_slide_2(prs)

    out = "WifipoolChatbot_Update.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
